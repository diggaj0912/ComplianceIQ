from pathlib import Path

from pptx import Presentation
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
OUTPUT = ROOT / "ComplianceIQ_Hackathon_Pitch.pptx"
HERO = ROOT / "backend" / "src" / "assets" / "hero.png"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

NAVY = RGBColor(15, 23, 42)
SLATE = RGBColor(51, 65, 85)
MUTED = RGBColor(100, 116, 139)
WHITE = RGBColor(255, 255, 255)
INDIGO = RGBColor(99, 102, 241)
VIOLET = RGBColor(139, 92, 246)
PURPLE = RGBColor(168, 85, 247)
GREEN = RGBColor(34, 197, 94)
AMBER = RGBColor(245, 158, 11)
RED = RGBColor(239, 68, 68)
BG = RGBColor(248, 250, 252)
LINE = RGBColor(226, 232, 240)


def set_bg(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_box(slide, left, top, width, height, fill, line=None, radius=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(1)
    return shape


def add_textbox(slide, left, top, width, height, text, size=24, bold=False, color=NAVY,
                align=PP_ALIGN.LEFT, font_name="Aptos", italic=False):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    f = r.font
    f.name = font_name
    f.size = Pt(size)
    f.bold = bold
    f.italic = italic
    f.color.rgb = color
    return box


def add_bullets(slide, left, top, width, height, items, size=18, color=SLATE, bullet_color=INDIGO):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.clear()
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.name = "Aptos"
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(8)
        p.bullet = True
    return box


def add_header(slide, title, section=None):
    set_bg(slide)
    add_box(slide, Inches(0.45), Inches(0.35), Inches(12.45), Inches(0.08), INDIGO)
    add_textbox(slide, Inches(0.55), Inches(0.55), Inches(8.5), Inches(0.6), title, size=28, bold=True, color=NAVY)
    if section:
        add_textbox(slide, Inches(10.0), Inches(0.52), Inches(2.5), Inches(0.4), section, size=12, bold=True, color=INDIGO, align=PP_ALIGN.RIGHT)


def add_metric_card(slide, left, top, width, height, label, value, sub=None, fill=WHITE, accent=INDIGO, value_size=26):
    add_box(slide, left, top, width, height, fill, LINE, radius=True)
    add_box(slide, left, top, Inches(0.1), height, accent, accent, radius=True)
    add_textbox(slide, left + Inches(0.22), top + Inches(0.18), width - Inches(0.3), Inches(0.3), label, size=12, bold=True, color=MUTED)
    add_textbox(slide, left + Inches(0.22), top + Inches(0.5), width - Inches(0.3), Inches(0.42), value, size=value_size, bold=True, color=NAVY)
    if sub:
        add_textbox(slide, left + Inches(0.22), top + Inches(1.0), width - Inches(0.3), Inches(0.35), sub, size=10, color=MUTED)


def title_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, NAVY)
    add_box(slide, Inches(0), Inches(0), prs.slide_width, prs.slide_height, NAVY)

    # Accent glow bands
    add_box(slide, Inches(0.7), Inches(0.9), Inches(4.2), Inches(0.12), VIOLET)
    add_box(slide, Inches(0.7), Inches(1.15), Inches(2.8), Inches(0.08), INDIGO)

    add_textbox(slide, Inches(0.7), Inches(1.45), Inches(6.0), Inches(0.7), "ComplianceIQ", size=34, bold=True, color=WHITE)
    add_textbox(slide, Inches(0.7), Inches(2.1), Inches(5.5), Inches(0.6), "AI-Powered GST Compliance Agent", size=24, bold=True, color=RGBColor(199, 210, 254))
    add_textbox(slide, Inches(0.7), Inches(2.75), Inches(5.9), Inches(1.0),
                "Turns raw transaction data into audit-ready GST filings in seconds.",
                size=18, color=RGBColor(226, 232, 240))

    # Stats strip
    stats = [
        ("40+ hrs/month", "saved for tax teams"),
        ("100%", "audit trail coverage"),
        ("₹1.8L", "avg ITC savings/cycle"),
    ]
    x = 0.7
    for value, sub in stats:
        add_box(slide, Inches(x), Inches(4.2), Inches(1.8), Inches(1.05), RGBColor(30, 41, 59), RGBColor(51, 65, 85), radius=True)
        add_textbox(slide, Inches(x + 0.12), Inches(4.38), Inches(1.55), Inches(0.32), value, size=18, bold=True, color=WHITE)
        add_textbox(slide, Inches(x + 0.12), Inches(4.7), Inches(1.55), Inches(0.26), sub, size=9, color=RGBColor(203, 213, 225))
        x += 2.0

    # Hero panel
    add_box(slide, Inches(8.0), Inches(0.8), Inches(4.55), Inches(5.9), RGBColor(17, 24, 39), RGBColor(71, 85, 105), radius=True)
    if HERO.exists():
        slide.shapes.add_picture(str(HERO), Inches(9.0), Inches(1.05), height=Inches(3.0))
    add_textbox(slide, Inches(8.35), Inches(4.25), Inches(3.8), Inches(0.3), "India-first. Offline-ready. Enterprise-grade.", size=14, bold=True, color=RGBColor(199, 210, 254), align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(8.35), Inches(4.65), Inches(3.8), Inches(0.75), "Built in 18 hours for hackathons, designed for real compliance teams.", size=16, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(8.35), Inches(5.55), Inches(3.8), Inches(0.35), "ComplianceIQ", size=20, bold=True, color=RGBColor(167, 139, 250), align=PP_ALIGN.CENTER)

    return slide


def problem_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "The Problem", "Why compliance breaks")

    add_metric_card(slide, Inches(0.6), Inches(1.3), Inches(3.85), Inches(1.8), "Annual Loss", "₹2.3 lakh crore", "Due to GST filing errors, missed ITC, and penalties", accent=RED)
    add_metric_card(slide, Inches(4.7), Inches(1.3), Inches(3.85), Inches(1.8), "Manual Effort", "40+ hrs/month", "Spent parsing CSVs and calculating liabilities", accent=AMBER)
    add_metric_card(slide, Inches(8.8), Inches(1.3), Inches(3.85), Inches(1.8), "Complexity", "5 GST slabs", "0 / 5 / 12 / 18 / 28% logic causes errors", accent=VIOLET)

    add_box(slide, Inches(0.6), Inches(3.45), Inches(12.1), Inches(2.45), WHITE, LINE, radius=True)
    add_textbox(slide, Inches(0.95), Inches(3.75), Inches(4.6), Inches(0.4), "What teams struggle with", size=18, bold=True, color=NAVY)
    add_bullets(slide, Inches(0.95), Inches(4.18), Inches(5.6), Inches(1.35), [
        "CSV data arrives messy and inconsistent",
        "GST slab categorization happens manually",
        "ITC claims are often incomplete or delayed",
        "Audits require evidence scattered across tools",
    ], size=16)

    add_box(slide, Inches(6.85), Inches(3.75), Inches(5.1), Inches(1.7), RGBColor(239, 246, 255), RGBColor(191, 219, 254), radius=True)
    add_textbox(slide, Inches(7.15), Inches(4.0), Inches(4.5), Inches(0.4), "Traditional compliance tools are static dashboards.", size=16, bold=True, color=NAVY)
    add_textbox(slide, Inches(7.15), Inches(4.42), Inches(4.5), Inches(0.55), "They show numbers. They do not explain anomalies, optimize ITC, or answer questions in real time.", size=14, color=SLATE)

    return slide


def solution_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Our Solution", "How ComplianceIQ works")

    steps = [
        ("1", "Upload CSV", "Bring in raw transaction data"),
        ("2", "Rules Engine", "Auto-categorize across GST slabs"),
        ("3", "AI ITC", "Optimize eligible tax credits"),
        ("4", "Ask Anything", "Conversational compliance answers"),
        ("5", "Audit Log", "Immutable activity trail"),
    ]
    lefts = [0.6, 3.0, 5.4, 7.8, 10.2]
    for i, (num, title, desc) in enumerate(steps):
        add_box(slide, Inches(lefts[i]), Inches(1.7), Inches(2.0), Inches(2.35), WHITE, LINE, radius=True)
        add_box(slide, Inches(lefts[i] + 0.65), Inches(1.92), Inches(0.7), Inches(0.7), INDIGO, INDIGO, radius=True)
        add_textbox(slide, Inches(lefts[i] + 0.82), Inches(2.08), Inches(0.35), Inches(0.25), num, size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(lefts[i] + 0.15), Inches(2.75), Inches(1.7), Inches(0.3), title, size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(lefts[i] + 0.12), Inches(3.05), Inches(1.75), Inches(0.6), desc, size=11, color=MUTED, align=PP_ALIGN.CENTER)
        if i < len(steps) - 1:
            add_textbox(slide, Inches(lefts[i] + 2.05), Inches(2.48), Inches(0.25), Inches(0.25), "→", size=22, bold=True, color=INDIGO, align=PP_ALIGN.CENTER)

    add_box(slide, Inches(0.8), Inches(4.6), Inches(11.7), Inches(1.35), RGBColor(248, 250, 252), LINE, radius=True)
    add_textbox(slide, Inches(1.1), Inches(4.87), Inches(11.0), Inches(0.5),
                "Result: audit-ready GST filings, faster decisions, and fewer compliance surprises.",
                size=20, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    return slide


def features_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Feature Set", "What the demo includes")

    cards = [
        ("Real-time dashboard", "KPI cards, risk score, GST trend charts", INDIGO),
        ("CSV upload", "Drag-and-drop import with auto-classification", VIOLET),
        ("GST Centre", "Liability, ITC, and rate bucket analysis", PURPLE),
        ("Compliance calendar", "Upcoming deadlines and due-date tracking", GREEN),
        ("AI compliance agent", "Conversational Q&A over your data", AMBER),
        ("Audit trail", "Immutable action log for every step", RED),
    ]
    positions = [(0.65, 1.35), (4.35, 1.35), (8.05, 1.35), (0.65, 3.45), (4.35, 3.45), (8.05, 3.45)]
    for (title, desc, accent), (x, y) in zip(cards, positions):
        add_box(slide, Inches(x), Inches(y), Inches(3.1), Inches(1.75), WHITE, LINE, radius=True)
        add_box(slide, Inches(x), Inches(y), Inches(0.12), Inches(1.75), accent, accent, radius=True)
        add_textbox(slide, Inches(x + 0.2), Inches(y + 0.2), Inches(2.6), Inches(0.25), title, size=15, bold=True, color=NAVY)
        add_textbox(slide, Inches(x + 0.2), Inches(y + 0.55), Inches(2.6), Inches(0.65), desc, size=11.5, color=SLATE)

    add_box(slide, Inches(0.7), Inches(5.65), Inches(11.9), Inches(0.9), RGBColor(239, 246, 255), RGBColor(191, 219, 254), radius=True)
    add_textbox(slide, Inches(0.95), Inches(5.9), Inches(11.4), Inches(0.3),
                "Everything is offline-first with mock fallback, so the demo stays alive even if the LLM or API is down.",
                size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    return slide


def architecture_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Architecture", "Offline-first and demo-safe")

    add_box(slide, Inches(0.75), Inches(1.4), Inches(3.0), Inches(4.8), WHITE, LINE, radius=True)
    add_textbox(slide, Inches(1.0), Inches(1.7), Inches(2.4), Inches(0.3), "Frontend", size=16, bold=True, color=INDIGO, align=PP_ALIGN.CENTER)
    add_bullets(slide, Inches(1.0), Inches(2.15), Inches(2.4), Inches(2.8), [
        "React 18 + Vite",
        "Tailwind CSS",
        "Recharts visualizations",
        "Conversational chat UI",
    ], size=15)

    add_box(slide, Inches(4.55), Inches(1.4), Inches(3.0), Inches(4.8), WHITE, LINE, radius=True)
    add_textbox(slide, Inches(4.8), Inches(1.7), Inches(2.4), Inches(0.3), "Backend", size=16, bold=True, color=VIOLET, align=PP_ALIGN.CENTER)
    add_bullets(slide, Inches(4.8), Inches(2.15), Inches(2.4), Inches(2.8), [
        "FastAPI",
        "SQLAlchemy + SQLite",
        "CSV ingestion",
        "Audit log persistence",
    ], size=15)

    add_box(slide, Inches(8.35), Inches(1.4), Inches(3.8), Inches(4.8), WHITE, LINE, radius=True)
    add_textbox(slide, Inches(8.65), Inches(1.7), Inches(3.2), Inches(0.3), "AI Layer", size=16, bold=True, color=PURPLE, align=PP_ALIGN.CENTER)
    add_bullets(slide, Inches(8.65), Inches(2.15), Inches(3.1), Inches(2.8), [
        "Groq LLaMA 3.3 70B",
        "Mock fallback for demos",
        "Routing-ready architecture",
        "Sub-second response target",
    ], size=15)

    add_textbox(slide, Inches(1.1), Inches(6.35), Inches(11.1), Inches(0.4),
                "Data flows from upload → rules engine → AI assistant → report/export, with every action written to an audit trail.",
                size=14, color=SLATE, align=PP_ALIGN.CENTER)
    return slide


def impact_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Impact & Metrics", "Why judges should care")

    metrics = [
        ("85%", "reduction in manual GST processing time", GREEN),
        ("100%", "audit trail coverage", INDIGO),
        ("<1 sec", "liability calculation across all slabs", VIOLET),
        ("₹1.8L", "average ITC savings per cycle", AMBER),
    ]
    pos = [(0.75, 1.4), (4.15, 1.4), (7.55, 1.4), (10.95, 1.4)]
    for (value, desc, accent), (x, y) in zip(metrics, pos):
        add_box(slide, Inches(x), Inches(y), Inches(2.15), Inches(2.0), WHITE, LINE, radius=True)
        add_textbox(slide, Inches(x + 0.15), Inches(y + 0.28), Inches(1.85), Inches(0.45), value, size=28, bold=True, color=accent, align=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(x + 0.1), Inches(y + 1.0), Inches(1.95), Inches(0.65), desc, size=11.5, color=SLATE, align=PP_ALIGN.CENTER)

    add_box(slide, Inches(0.8), Inches(4.0), Inches(12.0), Inches(2.0), RGBColor(15, 23, 42), RGBColor(30, 41, 59), radius=True)
    add_textbox(slide, Inches(1.15), Inches(4.45), Inches(11.3), Inches(0.5),
                "ComplianceIQ changes GST work from spreadsheet-heavy operations into a conversational, AI-assisted workflow.",
                size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1.35), Inches(5.0), Inches(10.9), Inches(0.5),
                "The result is faster filing, fewer errors, stronger auditability, and clear financial savings.",
                size=14, color=RGBColor(199, 210, 254), align=PP_ALIGN.CENTER)
    return slide


def differentiation_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "What Makes Us Different", "Why this wins hackathons")

    add_box(slide, Inches(0.7), Inches(1.35), Inches(5.9), Inches(4.8), WHITE, LINE, radius=True)
    add_textbox(slide, Inches(1.0), Inches(1.65), Inches(5.2), Inches(0.3), "Static tools", size=18, bold=True, color=RED)
    add_bullets(slide, Inches(1.0), Inches(2.1), Inches(4.9), Inches(2.8), [
        "Show charts but do not explain them",
        "Require manual navigation to find answers",
        "Break when the LLM or API is unavailable",
        "Leave audit evidence scattered",
    ], size=15)

    add_box(slide, Inches(6.95), Inches(1.35), Inches(5.7), Inches(4.8), RGBColor(239, 246, 255), RGBColor(191, 219, 254), radius=True)
    add_textbox(slide, Inches(7.25), Inches(1.65), Inches(5.0), Inches(0.3), "ComplianceIQ", size=18, bold=True, color=INDIGO)
    add_bullets(slide, Inches(7.25), Inches(2.1), Inches(4.9), Inches(2.8), [
        "You can ask natural-language compliance questions",
        "AI responds with context from your actual data",
        "Offline mock fallback keeps demos alive",
        "Every action is tracked for audit readiness",
    ], size=15)

    add_box(slide, Inches(2.25), Inches(5.55), Inches(8.8), Inches(0.95), WHITE, LINE, radius=True)
    add_textbox(slide, Inches(2.55), Inches(5.82), Inches(8.2), Inches(0.3),
                "Static dashboard + AI assistant + audit trail + offline demo safety = a memorable hackathon story.",
                size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    return slide


def roadmap_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Roadmap & Ask", "Close strong")

    add_box(slide, Inches(0.75), Inches(1.4), Inches(6.1), Inches(4.85), WHITE, LINE, radius=True)
    add_textbox(slide, Inches(1.05), Inches(1.7), Inches(5.5), Inches(0.3), "Immediate roadmap", size=18, bold=True, color=INDIGO)
    add_bullets(slide, Inches(1.05), Inches(2.15), Inches(5.2), Inches(3.2), [
        "Add live GST filings and export integrations",
        "Broaden ITC rules and anomaly detection",
        "Add multi-user roles for finance teams",
        "Connect to ERP and accounting systems",
    ], size=16)

    add_box(slide, Inches(7.15), Inches(1.4), Inches(5.45), Inches(4.85), RGBColor(15, 23, 42), RGBColor(30, 41, 59), radius=True)
    add_textbox(slide, Inches(7.45), Inches(1.75), Inches(4.9), Inches(0.3), "The ask", size=18, bold=True, color=RGBColor(199, 210, 254))
    add_textbox(slide, Inches(7.45), Inches(2.3), Inches(4.9), Inches(0.9),
                "Choose ComplianceIQ when you want a product that looks sharp, demos live, and tells a strong India-first AI story.",
                size=22, bold=True, color=WHITE)
    add_textbox(slide, Inches(7.45), Inches(3.65), Inches(4.9), Inches(0.6),
                "Built in 18 hours. Production-ready architecture. Hackathon-winning narrative.",
                size=14, color=RGBColor(199, 210, 254))
    add_box(slide, Inches(7.45), Inches(4.6), Inches(2.1), Inches(0.7), GREEN, GREEN, radius=True)
    add_textbox(slide, Inches(7.6), Inches(4.84), Inches(1.8), Inches(0.2), "Ready to demo", size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    return slide


slides = [title_slide, problem_slide, solution_slide, features_slide, architecture_slide, impact_slide, differentiation_slide, roadmap_slide]
for make_slide in slides:
    make_slide()

prs.save(OUTPUT)
print(f"Saved {OUTPUT}")